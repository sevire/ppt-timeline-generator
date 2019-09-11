#!/usr/bin/env python
# coding: utf-8

import os
from pathlib import Path

from get_command_line_parameters import get_command_line_parameters
from milestone_utilities import create_milestone_shapes
from pptx_utilities import extract_template_data
from pptx import Presentation
import logging
import time
# Not sure where to import MSO_SHAPE or MSO_CONNECTOR from
import read_data_utilities as read_data
import sys

timeline_param = {}


def main():
    ts = time.gmtime()
    time_string = time.strftime("%Y_%m_%d__%H_%M_%S", ts)
    logging.basicConfig(filename='resources/logs/create_milestone_{}.log'.format(time_string), level=logging.DEBUG)

    root_dir_old = Path(
        '/Users/thomasgaylardou/OneDrive - The Open University/'
        'BusinessDesign/Vision/MediumLongTerm/ACLA/EPA/timeline/production')

    root_dir = Path('/Users/thomasgaylardou/Documents/EPA-timeline/mags-lois-testing/production')

    wb_path, template_dir, output_dir = get_command_line_parameters()

    # workbook_dir = root_dir
    # workbook_name = 'ApprenticeshipTimelineData-v01-CMDA.xlsx'

    template_dir = root_dir / 'templates'
    output_dir = root_dir / 'timelines'

    xl_object, timeline_list = read_data.get_list_of_timeline_sheets(wb_path)

    for timeline_id in timeline_list:
        parameters = read_data.read_parameter_data(xl_object, timeline_id)
        read_data.extract_parameter_data(timeline_id, timeline_param, parameters)

        milestone_data = read_data.read_milestone_data(xl_object, timeline_id)

        presentation_name_in = os.path.join(template_dir, timeline_id + '.pptx')
        presentation_name_out = os.path.join(output_dir, timeline_id + '-out.pptx')
        prs = Presentation(presentation_name_in)

        slides = prs.slides

        # Expect to find template objects in first slide - don't look anywhere else
        template_slide = slides[0]
        template_shapes = template_slide.shapes
        template_data = extract_template_data(template_shapes)

        # We are adding shapes to the second slide.
        milestone_slide = slides[1]
        milestone_shapes = milestone_slide.shapes
        create_milestone_shapes(timeline_id, milestone_data, milestone_shapes, template_data, timeline_param)

        prs.save(presentation_name_out)


if __name__ == '__main__':
    main()
