#!/usr/bin/env python
# coding: utf-8

import os
from pathlib import Path

from create_timeline.get_command_line_parameters import get_command_line_parameters
from create_timeline.milestone_utilities import create_milestone_shapes
from create_timeline.pptx_utilities import extract_template_data
from pptx import Presentation
import logging
import time
# Not sure where to import MSO_SHAPE or MSO_CONNECTOR from
from create_timeline import read_data_utilities as read_data

timeline_param = {}
logger = logging.getLogger()


def main():
    ts = time.gmtime()
    time_string = time.strftime("%Y_%m_%d__%H_%M_%S", ts)
    logging.basicConfig(filename='resources/logs/create_milestone.log'.format(time_string),
                        filemode='w',
                        format='%(levelname)s: %(asctime)s %(message)s',
                        level=logging.DEBUG
                        )

    wb_path, template_dir, output_dir = get_command_line_parameters()

    logger.info('EPA Timeline Generator starting...\n\n')
    logger.info(f'Excel file is:      {wb_path}')
    logger.info(f'Template folder is: {template_dir}')
    logger.info(f'Output folder is:   {output_dir}\n')

    logger.info(f'Reading Excel driver file')
    xl_object, timeline_list = read_data.get_list_of_timeline_sheets(wb_path)

    logger.info('Timeline driver data found for following timeline names...')
    logger.info(timeline_list)

    for timeline_id in timeline_list:
        logger.info('\n')
        logger.info(f'Processing timeline {timeline_id}.')
        parameters = read_data.read_parameter_data(xl_object, timeline_id)
        read_data.extract_parameter_data(timeline_id, timeline_param, parameters)

        milestone_data = read_data.read_milestone_data(xl_object, timeline_id)

        presentation_name_in = os.path.join(template_dir, timeline_id + '.pptx')
        presentation_name_out = os.path.join(output_dir, timeline_id + '-out.pptx')
        logger.info(f'Template file for {timeline_id} is {presentation_name_in}')
        logger.info(f'Output file for {timeline_id} is {presentation_name_out}')
        prs = Presentation(presentation_name_in)

        slides = prs.slides

        # Expect to find template objects in first slide - don't look anywhere else
        template_slide = slides[0]
        template_shapes = template_slide.shapes
        template_data = extract_template_data(template_shapes)

        # We are adding shapes to the second slide.
        milestone_slide = slides[1]
        milestone_shapes = milestone_slide.shapes
        create_milestone_shapes(timeline_id, milestone_data, milestone_shapes, template_data, timeline_param)

        prs.save(presentation_name_out)


if __name__ == '__main__':
    main()
