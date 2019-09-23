from pptx.dml.color import RGBColor, _NoneColor
from pptx.util import Pt
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
import logging

logger = logging.getLogger()


def extract_shape_attributes(src_shape):
    """
    Takes an existing (template) autoshape and extracts formatting properties which can later be applied to
    another shape (with modifications if required).

    Note this is very simplified by making assumptions about what is likely to be needed. So for example, the fill
    colour is assumed to always be a solid colour using RGB values (for now).
    """

    fill = src_shape.fill
    line = src_shape.line
    fore_colour = fill.fore_color
    line_colour = src_shape.line.color
    logger.debug('line_colour = {}'.format(line_colour))
    if isinstance(line_colour._color, _NoneColor):
        line_colour_rgb = None
    else:
        line_colour_rgb = line_colour.rgb
    run = src_shape.text_frame.paragraphs[0].runs[0]

    shape_attributes = {
        'shape_type': src_shape.shape_type,
        'autoshape_type': src_shape.auto_shape_type,
        'rotation': src_shape.rotation,
        'left': src_shape.left,
        'top': src_shape.top,
        'width': src_shape.width,
        'height': src_shape.height,
        'fill_colour_rgb': fill.fore_color.rgb,
        'line_colour_rgb': line_colour_rgb,
        # 'line_colour_brightness': line.color.brightness,
        'line_width': line.width,
        'font_name': run.font.name,
        'font_size': run.font.size,
        'font_bold': run.font.bold,
        'font_italic': run.font.italic,
        'font_colour_rgb': run.font.color.rgb
    }

    return shape_attributes


# def manage_shape_line_attributes(line_object):
#     """
#     Takes the line attribute of a shape (e.g. rectangle) and manages setting and reading of
#
#     :param attribute_object:
#     :return:
#     """

def apply_attribute(shape_object, shape_attributes, attribute_name):
    if attribute_name in shape_attributes and shape_attributes[attribute_name] is not None:
        setattr(shape_object, attribute_name, shape_attributes[attribute_name])
    else:
        setattr(shape_object, attribute_name, None)


def apply_shape_attributes(target_shape, shape_attributes, include_construction_attrs=False):
    """
    Takes a supplied (typically newly created shape) and a dictionary of attributes and applies the attributes
    (or those which have a value) to the supplied shape.

    It's all done pretty mindlessly with each attribute processed individually.
    """
    # logger.debug('func:apply_shape_attributes, shape: {}, attributes: {}'.format(target_shape, shape_attributes))

    # Usually don't need width etc as that would have been set at creation of shape.
    if include_construction_attrs:
        apply_attribute(target_shape, shape_attributes, 'width')
        apply_attribute(target_shape, shape_attributes, 'left')
        apply_attribute(target_shape, shape_attributes, 'top')
        apply_attribute(target_shape, shape_attributes, 'height')

        # if 'width' in shape_attributes and shape_attributes['width'] is not None:
        #     target_shape.width = shape_attributes['width']
        # if 'left' in shape_attributes and shape_attributes['left'] is not None:
        #     target_shape.width = shape_attributes['left']
        #
        # if 'top' in shape_attributes and shape_attributes['top'] is not None:
        #     target_shape.width = shape_attributes['top']
        #
        # if 'height' in shape_attributes and shape_attributes['height'] is not None:
        #     target_shape.width = shape_attributes['height']

    if 'rotation' in shape_attributes and shape_attributes['rotation'] is not None:
        target_shape.rotation = shape_attributes['rotation']

    if 'fill_colour_rgb' in shape_attributes and shape_attributes['fill_colour_rgb'] is not None:
        fill = target_shape.fill
        fill.solid()
        fore_color = fill.fore_color
        fore_color.rgb = shape_attributes['fill_colour_rgb']

    if 'line_colour_rgb' in shape_attributes:
        line = target_shape.line
        if shape_attributes['line_colour_rgb'] is None:
            line.fill.background()
        else:
            line.color.rgb = shape_attributes['line_colour_rgb']

    if 'line_colour_brightness' in shape_attributes and shape_attributes['line_colour_brightness'] is not None:
        line = target_shape.line
        line.color.brightness = shape_attributes['line_colour_brightness']

    if 'line_width' in shape_attributes and shape_attributes['line_width'] is not None:
        line = target_shape.line
        line.width = shape_attributes['line_width']

    # Note: Only handles shapes with one run.
    run = target_shape.text_frame.paragraphs[0].runs[0]

    if 'font_name' in shape_attributes and shape_attributes['font_name'] is not None:
        run.font.name = shape_attributes['font_name']

    if 'font_size' in shape_attributes and shape_attributes['font_size'] is not None:
        run.font.size = shape_attributes['font_size']

    if 'font_bold' in shape_attributes and shape_attributes['font_bold'] is not None:
        run.font.bold = shape_attributes['font_bold']

    if 'font_italic' in shape_attributes and shape_attributes['font_italic'] is not None:
        run.font.italic = shape_attributes['font_italic']

    if 'font_colour_rgb' in shape_attributes and shape_attributes['font_colour_rgb'] is not None:
        run.font.color.rgb = shape_attributes['font_colour_rgb']


def create_shape_from_template(shapes_object, left, top, shape_attributes, text):
    """
    Adds an autoshape to a supplied python-pptx shapes object.  The attributes of the shape are supplied and
    will have been previously extracted from a template shape which sets the formatting for all similar shapes.

    The top and left attributes are supplied separately as each new shape will have its own position.

    """
    arguments = (shape_attributes['autoshape_type'],
                 left,
                 top,
                 shape_attributes['width'],
                 shape_attributes['height'],
                 )
    new_shape = shapes_object.add_shape(*arguments)

    # Add text here, not in attributes
    text_frame = new_shape.text_frame
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.line_spacing = 0.8
    run = p.add_run()
    run.text = text

    # Hard code text formatting for now until we have implemented capturing these attributes
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(8)
    font.bold = True
    font.italic = False  # cause value to be inherited from theme
    font.color.rgb = RGBColor(0xFF, 0, 0)

    # Now apply remainder of formatting to shape
    apply_shape_attributes(new_shape, shape_attributes)


def create_line(shapes_object, x, y_ms, y_text, colour):
    """
    Creates a vertical line between a milestone and it's associated text box.

    There are two points to define but as the line will always be vertical the x coordinate is the same for both.

    :param shapes_object:
    :param x: x coordinate which will be the centre of the milestone
    :param y_ms: y coordinate of the end which joins the text box
    :param y_text: y coordinate of the end which joins the milestone
    :param colour:
    :return:
    """

    line = shapes_object.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x, y_ms, x, y_text)
    line.line.color.rgb = RGBColor(0, 32, 96)
    line.line.color.rgb = colour


def extract_template_data(shapes_object):
    expected_templates = {'1': {}, '2': {}, '3': {}, 'TEXT 1': {}, 'TEXT 2': {}, 'TEXT 3': {}}

    for shape_number, template_shape in enumerate(shapes_object):
        logger.debug('{:3}: Shape type: {}'.format(shape_number, template_shape.shape_type))

        if template_shape.has_text_frame:
            shape_text = template_shape.text

            # Now work out which template shape it is and extract and store the key formatting information
            # Note: I am taking this approach as at the time of writing there is no way of simply copying a shape
            # using python-pptx
            if shape_text in expected_templates:
                if 'data' in expected_templates[shape_text]:
                    logger.warning('Duplicate template found for {}, ignoring'.format(shape_text))
                else:
                    # We've found a valid template so store it's data
                    logger.info('Found milestone shape {}'.format(shape_text))
                    expected_templates[shape_text] = {}
                    expected_templates[shape_text]['data'] = extract_shape_attributes(template_shape)
            else:
                 logger.info('Not a template')
        else:
            logger.info('No text frame')

    # We should have all templates now, just tidy up a bit and tell the user
    undiscovered = list(filter(lambda x: 'data' not in expected_templates[x], expected_templates))
    if len(undiscovered) > 0:
        logger.warning('templates not found: {}'.format(undiscovered))
    logger.info(expected_templates)

    return expected_templates
