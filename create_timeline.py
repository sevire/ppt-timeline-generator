#!/usr/bin/env python
# coding: utf-8

import os
from milestone_utilities import create_milestone_shapes
from pptx_utilities import extract_template_data
from pptx import Presentation
import logging
import time

# Not sure where to import MSO_SHAPE or MSO_CONNECTOR from
from read_data_utilities import read_milestone_data


def main():
    ts = time.gmtime()
    time_string = time.strftime("%Y-%m-%d_%H:%M:%S", ts)
    logging.basicConfig(filename='resources/logs/create_milestone_{}.log'.format(time_string), level=logging.DEBUG)

    root_dir_old = \
        '/Users/thomasgaylardou/OneDrive - The Open University/' \
        'BusinessDesign/Vision/MediumLongTerm/ACLA/EPA/timeline/production'

    root_dir = \
        '/Users/thomasgaylardou/Documents/EPA-timeline/production'

    workbook_dir = root_dir
    workbook_name = 'ApprenticeshipTimelineData-v01.xlsx'

    template_dir = os.path.join(root_dir, 'templates')
    output_dir = os.path.join(root_dir, 'timelines')

    for timeline_id in ['CPT-HAHAP-01', 'CPT-HAHAP-02']:

        milestone_excel_workbook_name = os.path.join(workbook_dir, workbook_name)
        milestones_data = read_milestone_data(milestone_excel_workbook_name, sheet_name=timeline_id)

        presentation_name_in = os.path.join(template_dir, timeline_id + '.pptx')
        presentation_name_out = os.path.join(output_dir, timeline_id + '-out.pptx')
        prs = Presentation(presentation_name_in)

        slides = prs.slides

        # Expect to find template objects in first slide - don't look anywhere else
        template_slide = slides[0]
        template_shapes = template_slide.shapes
        template_data = extract_template_data(template_shapes)

        milestone_slide = slides[1]
        milestone_shapes = milestone_slide.shapes
        create_milestone_shapes(timeline_id, milestones_data, milestone_shapes, template_data)

        prs.save(presentation_name_out)


if __name__ == '__main__':
    main()
