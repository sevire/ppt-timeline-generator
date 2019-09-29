from pptx.dml.color import RGBColor, _NoneColor
from datetime import datetime as dt
import logging

from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.util import Pt

logger = logging.getLogger(__name__)


class PptTimelineGeneratedSlide:
    def __init__(self, slide, timeline):
        self.slide = slide
        self.timeline = timeline

    def generate_plot_data(self, templates):
        """
        Iterate through the milestone data and calculate the information required to create new shapes.  Store it until
        all the data has been collected, then go back through the data three times to create the shapes in the order of:
        - Connectors
        - Milestones
        - Text boxes

        As new shapes are placed in front of older shapes, this will ensure that the connectors travel behind the shapes.

        :return: None
        """

        # track_orientation = 1

        # Create the structure to store all the data to plot the points, but plot later to control what appears on top
        # of what.
        shape_data = {
            'connector': {
                'func': self.create_line,
                'parameters': []
            },
            'milestone': {
                'func': self.create_shape_from_template,
                'parameters': []
            },
            'textbox': {
                'func': self.create_shape_from_template,
                'parameters': []
            }
        }
        for milestone_number, milestone_data in self.timeline.milestone_data.iterrows():
            milestone_text = milestone_data['Milestone Name']
            milestone_date = milestone_data['Date']
            milestone_level = milestone_data['Milestone Level']  # Normalises track to centre on zero
            milestone_level_text_label = 'TEXT {}'.format(milestone_level)
            milestone_track_override = self.parse_track_override(milestone_data['Milestone Track Override'])
            textbox_track_override = self.parse_track_override(milestone_data['Textbox Track Override'])

            # Plot milestone shape at right position
            day_num = self.to_day_num(milestone_date)

            # Calculate milestone position and whether above or below centre line
            milestone_track_number = milestone_track_override if milestone_track_override is not None \
                else (milestone_level - 1) * self.timeline.parameter_data['milestone_track_orientation'][milestone_level]

            milestone_data, textbox_data = \
                self.get_plot_data(
                    day_num, milestone_track_number,
                    self.timeline.parameter_data['milestone_track_orientation'][milestone_level],
                    textbox_track_override
                )

            # Create and position milestone marker shape
            ms_x, ms_y = milestone_data

            milestone_width = templates[str(milestone_level)]['data']['width']
            milestone_height = templates[str(milestone_level)]['data']['height']

            # Extract fill color of milestone as we will re-use this for connector (for the moment)
            milestone_colour = templates[str(milestone_level)]['data']['fill_colour_rgb']

            ms_left = ms_x - milestone_width // 2
            ms_top = ms_y - milestone_height // 2

            milestone_contents = self.generate_milestone_text(milestone_number)

            # Add milestone data to data structure for plotting later
            shape_data['milestone']['parameters'].append(
                (self.slide.shapes, ms_left, ms_top, templates[str(milestone_level)]['data'], milestone_contents)
            )

            # Create and position text box shape
            text_x, text_y = textbox_data

            textbox_width = templates[milestone_level_text_label]['data']['width']
            textbox_height = templates[milestone_level_text_label]['data']['height']

            textbox_shift = self.calculate_text_box_shift(
                                                     day_num,
                                                     textbox_width,
                                                     milestone_width
                                                     )

            tx_left = text_x - textbox_width // 2 + textbox_shift
            tx_top = text_y - textbox_height // 2

            display_date = dt.strftime(milestone_date, '%d-%b-%Y')
            textbox_text = self.generate_textbox_text(
                display_date,
                milestone_text,
                milestone_number,
                False,
                daynum=day_num,
                daynum_prop=self.day_num_as_proportion(day_num),
                shift=textbox_shift)

            # Add text box data to data structure for plotting later
            shape_data['textbox']['parameters'].append(
                (self.slide.shapes, tx_left, tx_top, templates['TEXT {}'.format(milestone_level)]['data'], textbox_text)
            )

            # Create and position connecting line.  It won't actually connect to anything but will be placed correctly
            line_x = ms_x

            # Plot connector from middle of shapes (for now) to avoid having to re-factor to properly process tracks
            # ToDo: Refactor this so that tracks of ms and tx are used to calculate connector line

            line_ms_y = ms_y
            line_text_y = text_y

            # Add connector data to data structure for plotting later
            shape_data['connector']['parameters'].append(
                (self.slide.shapes, line_x, line_ms_y, line_text_y, milestone_colour))

            # Alternate between positive and negative track numbers to minimise overlap of objects
            self.timeline.parameter_data['milestone_track_orientation'][milestone_level] = \
                -self.timeline.parameter_data['milestone_track_orientation'][milestone_level]

        # Now plot the data in the order of connectors (bottom), milestones and text boxes
        for shape_type in shape_data:
            func = shape_data[shape_type]['func']
            for param in shape_data[shape_type]['parameters']:
                func(*param)

    @staticmethod
    def create_line(shapes_object, x, y_ms, y_text, colour):
        """
        Creates a vertical line between a milestone and it's associated text box.

        There are two points to define but as the line will always be vertical the x coordinate is the same for both.

        :param shapes_object:
        :param x: x coordinate which will be the centre of the milestone
        :param y_ms: y coordinate of the end which joins the text box
        :param y_text: y coordinate of the end which joins the milestone
        :param colour:
        :return:
        """

        line = shapes_object.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x, y_ms, x, y_text)
        line.line.color.rgb = RGBColor(0, 32, 96)
        line.line.color.rgb = colour

    def create_shape_from_template(self, shapes_object, left, top, shape_attributes, text):
        """
        Adds an autoshape to a supplied python-pptx shapes object.  The attributes of the shape are supplied and
        will have been previously extracted from a template shape which sets the formatting for all similar shapes.

        The top and left attributes are supplied separately as each new shape will have its own position.

        """
        arguments = (shape_attributes['autoshape_type'],
                     left,
                     top,
                     shape_attributes['width'],
                     shape_attributes['height'],
                     )
        new_shape = shapes_object.add_shape(*arguments)

        # Add text here, not in attributes
        text_frame = new_shape.text_frame
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.line_spacing = 0.8
        run = p.add_run()
        run.text = text

        # Hard code text formatting for now until we have implemented capturing these attributes
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(8)
        font.bold = True
        font.italic = False  # cause value to be inherited from theme
        font.color.rgb = RGBColor(0xFF, 0, 0)

        # Now apply remainder of formatting to shape
        self.apply_shape_attributes(new_shape, shape_attributes)

    @staticmethod
    def parse_track_override(override_value):
        """
        The spreadsheet can contain a track number to override the calculated one for either milestone or text box.  This
        can help with fine tuning positioning for overlapping elements or just to control which elements line up with
        which others.

        If the override cell is left blank in a particular row, then it will come through as Nan.  We need to convert to
        None.  Also we can check that the track number isn't a stupid value.

        :param override_value:
        :return: Integer track number or None
        """
        try:
            track_number = int(override_value)
            return track_number
        except ValueError:
            return None
        except Exception:
            print('Unexpected exception converting track_override {}'.format(override_value))
            return None

    def generate_milestone_text(self, milestone_number):
        if self.timeline.parameter_data['include_ms_num_in_ms'] is True:
            return f'{milestone_number}'
        else:
            return ''

    def generate_textbox_text(self,
                              milestone_date,
                              milestone_text,
                              milestone_number,
                              testing_flag=False,
                              **kwargs
                              ):
        """
        Generates text for textbox.  Required as there are variations depending upon what options are configured.

        :param milestone_date:
        :param milestone_text:
        :param milestone_number: Number of current milestone being processed.  Optional inclusion.
        :param timeline_name: Name of timeline.  Included if not in debug mode.
        :param tlp: Passed in to get config data for this timeline
        :param testing_flag: If True then populate with data useful for testing, not actual data.
        :return:
        """

        if testing_flag is not True:
            if self.timeline.parameter_data['include_ms_num_in_text'] is True:
                return f'[{milestone_number}]-[{milestone_date}]\n{milestone_text}'
            else:
                return f'[{milestone_date}]\n{milestone_text}'
        else:
            # Test mode: Whatever data is passed in kwargs, include in text box
            return f'[{milestone_number}]-[{milestone_date}]\n{kwargs}'

    # ### Utility functions for working out where to place things
    def day_num_as_proportion(self, day):
        return (day - 1) / (self.timeline.parameter_data['milestone_total_days'] - 1)

    def to_day_num(self, date):
        start_date = self.timeline.parameter_data['start_date']
        return (date - start_date).days + 1

    @staticmethod
    def sign(x):
        return (1, -1)[x <= 0]  # Implement sign function as isn't standard

    def get_plot_data(self, day_num, milestone_track, orientation, textbox_track_override):
        """
        Calculates the position of each of the elements of the milestone, which are:
        - Milestone indicator (circle)
        - Milestone description (rectangle)
        - Connecting line (line)

        The arrangement needs to be as follows:

        The track number is a positive or negative integer which describes how far away the milestone element needs
        to be from the centre line, and whether it is above or below it.

        The milestone indicator will be positioned around the centre line at a vertical distance determined by the
        track number allocated to the milestone and a horizontal distance which places it precisely on the right date
        for the milestone.

        The text box will be placed in the same configuration as the milestone indicator, but further away from the
        centre line.

        The connecting line will run from the edge of the milestone indicator in the direction of its corresponding
        text box, to the nearest edge of the text box.

        What will be returned will be three tuples, which indicate where the CENTRE of the corresponding element
        needs to be.  The caller will then need to adjust depending upon the width and height of each element.

        The tuples are structured as follows:

        - milestone_indicator_data  = (left, top)
        - milestone_description_box = (left, top)
        - connecting_line           = (left_1, top_1, left_2, top_2)

        :param day_num:
        :param milestone_track:
        :param orientation:
        :param textbox_track_override:
        :return:
        """
        # Calculate x position first as this will be used in plotting of all elements.
        x_rel = self.day_num_as_proportion(day_num) * self.timeline.parameter_data['milestone_x_range']
        x_abs = int(self.timeline.parameter_data['milestone_left'] + x_rel)

        # Calculate data to drive plotting of each element

        # milestone_vertical_position = calculate_offset_for_ms_track(milestone_track)
        milestone_vertical_position = self.calculate_track_location(
            milestone_track,
            *self.timeline.parameter_data['milestone_track_data']
        )
        milestone_element_data = (x_abs, milestone_vertical_position)

        # --- Second the text box
        textbox_vertical_position = self.calculate_offset_for_text_track(orientation, day_num, textbox_track_override)
        textbox_element_data = (x_abs, textbox_vertical_position)

        return milestone_element_data, textbox_element_data  # Not including connector line yet

    def calculate_next_track(self, day_num, orientation):
        """
        Works out which part (third probably) of chart this day appears on then retrieves track value to use and
        updates for next call.

        :param day_num:
        :param orientation:
        :return:
        """
        area_num = self.chart_area_num(day_num)
        current_track = self.timeline.parameter_data['textbox_track_position_data'][area_num]['track'][orientation]
        self.timeline.parameter_data['textbox_track_position_data'][area_num]['track'][orientation] = \
            ((current_track - 1 + self.timeline.parameter_data['textbox_track_position_data'][area_num]['direction']) %
             self.timeline.parameter_data['num_text_tracks']) + 1

        return current_track * orientation

    def chart_area_num(self, day_num):
        """
        Calculates which part of the chart this day will appear in.  One of:
            1: Left
            2: Middle
            3: Right

        :param day_num: Which day within the date range for the chart this is
        :return: 1, 2, or 3.
        """
        return int(self.day_num_as_proportion(day_num) * 3 + 1)

    def calculate_offset_for_text_track(self, orientation, day_num, track_number=None):
        """
        Similar to calculate_offset_for_ms_track, although there is no level to base the track number on.  Instead
        we will choose the track using two factors:

        - The orientation (+1 or -1) will be passed in and the selected track number will conform to that so that the
          text box is on the same side of the centre line as the milestone.

        - A record will be kept of the most recent track number for each orientation and at each call the next number
          will be selected (mod the number of tracks per side of the centre line)

        :param orientation:
        :param day_num:
        :param track_number:
        :return:
        """
        if track_number is not None:
            track = track_number
        else:
            track = self.calculate_next_track(day_num, orientation)

        position = self.calculate_track_location(track, *self.timeline.parameter_data['textbox_track_data'])

        # Calculate the (relative) amount off the centre line we need to place the centre of the text box
        # vertical_offset_rel = text_box_track_start_offset + (track - 1) * offset_per_text_track

        # Now calculate the absolute offset based on where the centre line is and which orientation we are plotting
        # vertical_position = centre_vertical_position + orientation * vertical_offset_rel

        return position

    def calculate_track_location(self, track_number, track_separation, centre_position, track_start_offset=0):
        """
        General purpose function to calculate the vertical position of where the centre of an element needs to be to line
        up with the specified track.  Will be used for both milestone and textbox tracks (and others if required later).

        The way tracks will work is that different element types (e.g. milestones or associated text boxes) will be aligned
        to pre-determined horizontal 'tracks'.  Tracks can be above or below a specified centre line.

        If the tracks include the centre line then the centre line has track number zero.  For other cases a track of zero
        has no meaning and is invalid.

        If the tracks don't include the centre line then the first track will start at a specified distance from the centre
        line and all other tracks positioned relative to it.

        :param track_number:       Signed integer denoting the number of the track at which the element is to be placed.
                                   - 0  is the centre track but this is only meaningful when the set of tracks includes the
                                        centre (i.e if the track_start_offset is zero).
                                   - +n is the nth track below the centre line (using the PowerPoint plotting convention that
                                        the y coordinate starts at the top and increases as it goes down the page.
                                   - -n is the nth track above the centre line.

        :param track_separation:   The vertical distance between tracks

        :param centre_position:    Vertical position (supplied as pptx.util.Cm probably) of centre line of tracks.
                                   Used as anchor point to calculate the offsets for tracks from.

        :param track_start_offset: If the set of tracks doesn't include the centre line then this value specifies how far
                                   from the centre line the first track should be.

        :return:                   Distance from the top of the slide where the track is located.
        """
        orientation = self.sign(track_number)
        abs_track_num = abs(track_number)
        abs_track_offset = track_start_offset + (abs_track_num - 1) * track_separation
        position = centre_position + orientation * abs_track_offset

        return position

    def calculate_text_box_shift(self, day_num, text_box_width, milestone_width):
        """
        Calculates the correction to make to a text box's left position to take account of where it is on the slide.

        (NOTE: This has been updated to a continuous function rather than a step function based on which third
        of the screen the milestone is in.  This is to avoid the situation where two textboxes close to the
        threshold third line but on either side can be placed incorrectly in relation to each other.

        Additionally, it's important that the textbox isn't shifted so far that it isn't in line with the
        milestone marker - at each end of the timeline, the textbox edge should line up with the milestone
        edge.

        So this means that, at the left end of the timeline, the shift is:

          -->     +(textbox_width/2) - (ms_width/2)

        and at the right end of the timeline, the shift is:

          -->     -(textbox_width/2) + (ms_width/2)

        in between it's just a linear proportion between the two.

        :param day_num:         Used to calculate which third of the chart the text box is
        :param text_box_width:  Used to calculate how far to move the text box
        :param milestone_width: Used to help line up text box with edge of milestone shape
        :return:
        """

        shift_at_left = +(text_box_width / 2) - (milestone_width / 2)
        shift_at_right = -(text_box_width / 2) + (milestone_width / 2)

        shift_range = shift_at_right - shift_at_left
        proportion = self.day_num_as_proportion(day_num)

        shift = shift_at_left + proportion * shift_range

        return shift

    def apply_shape_attributes(self, target_shape, shape_attributes, include_construction_attrs=False):

        # (REFACTORED)

        """
        Takes a supplied (typically newly created shape) and a dictionary of attributes and applies the attributes
        (or those which have a value) to the supplied shape.

        It's all done pretty mindlessly with each attribute processed individually.
        """
        # logger.debug('func:apply_shape_attributes, shape: {}, attributes: {}'.format(target_shape, shape_attributes))

        # Usually don't need width etc as that would have been set at creation of shape.
        if include_construction_attrs:
            self.apply_attribute(target_shape, shape_attributes, 'width')
            self.apply_attribute(target_shape, shape_attributes, 'left')
            self.apply_attribute(target_shape, shape_attributes, 'top')
            self.apply_attribute(target_shape, shape_attributes, 'height')

            # if 'width' in shape_attributes and shape_attributes['width'] is not None:
            #     target_shape.width = shape_attributes['width']
            # if 'left' in shape_attributes and shape_attributes['left'] is not None:
            #     target_shape.width = shape_attributes['left']
            #
            # if 'top' in shape_attributes and shape_attributes['top'] is not None:
            #     target_shape.width = shape_attributes['top']
            #
            # if 'height' in shape_attributes and shape_attributes['height'] is not None:
            #     target_shape.width = shape_attributes['height']

        if 'rotation' in shape_attributes and shape_attributes['rotation'] is not None:
            target_shape.rotation = shape_attributes['rotation']

        if 'fill_colour_rgb' in shape_attributes and shape_attributes['fill_colour_rgb'] is not None:
            fill = target_shape.fill
            fill.solid()
            fore_color = fill.fore_color
            fore_color.rgb = shape_attributes['fill_colour_rgb']

        if 'line_colour_rgb' in shape_attributes:
            line = target_shape.line
            if shape_attributes['line_colour_rgb'] is None:
                line.fill.background()
            else:
                line.color.rgb = shape_attributes['line_colour_rgb']

        if 'line_colour_brightness' in shape_attributes and shape_attributes['line_colour_brightness'] is not None:
            line = target_shape.line
            line.color.brightness = shape_attributes['line_colour_brightness']

        if 'line_width' in shape_attributes and shape_attributes['line_width'] is not None:
            line = target_shape.line
            line.width = shape_attributes['line_width']

        # Note: Only handles shapes with one run.
        run = target_shape.text_frame.paragraphs[0].runs[0]

        if 'font_name' in shape_attributes and shape_attributes['font_name'] is not None:
            run.font.name = shape_attributes['font_name']

        if 'font_size' in shape_attributes and shape_attributes['font_size'] is not None:
            run.font.size = shape_attributes['font_size']

        if 'font_bold' in shape_attributes and shape_attributes['font_bold'] is not None:
            run.font.bold = shape_attributes['font_bold']

        if 'font_italic' in shape_attributes and shape_attributes['font_italic'] is not None:
            run.font.italic = shape_attributes['font_italic']

        if 'font_colour_rgb' in shape_attributes and shape_attributes['font_colour_rgb'] is not None:
            run.font.color.rgb = shape_attributes['font_colour_rgb']

    def extract_shape_attributes(self, src_shape):

        # (REFACTORED)

        """
        Takes an existing (template) autoshape and extracts formatting properties which can later be applied to
        another shape (with modifications if required).

        Note this is very simplified by making assumptions about what is likely to be needed. So for example, the fill
        colour is assumed to always be a solid colour using RGB values (for now).
        """

        fill = src_shape.fill
        line = src_shape.line
        fore_colour = fill.fore_color
        line_colour = src_shape.line.color
        logger.debug('line_colour = {}'.format(line_colour))
        if isinstance(line_colour._color, _NoneColor):
            line_colour_rgb = None
        else:
            line_colour_rgb = line_colour.rgb
        run = src_shape.text_frame.paragraphs[0].runs[0]

        shape_attributes = {
            'shape_type': src_shape.shape_type,
            'autoshape_type': src_shape.auto_shape_type,
            'rotation': src_shape.rotation,
            'left': src_shape.left,
            'top': src_shape.top,
            'width': src_shape.width,
            'height': src_shape.height,
            'fill_colour_rgb': fill.fore_color.rgb,
            'line_colour_rgb': line_colour_rgb,
            # 'line_colour_brightness': line.color.brightness,
            'line_width': line.width,
            'font_name': run.font.name,
            'font_size': run.font.size,
            'font_bold': run.font.bold,
            'font_italic': run.font.italic,
            'font_colour_rgb': run.font.color.rgb
        }

        return shape_attributes

    @staticmethod
    def apply_attribute(shape_object, shape_attributes, attribute_name):
        if attribute_name in shape_attributes and shape_attributes[attribute_name] is not None:
            setattr(shape_object, attribute_name, shape_attributes[attribute_name])
        else:
            setattr(shape_object, attribute_name, None)
