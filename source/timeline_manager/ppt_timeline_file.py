from pptx import Presentation
import logging

from source.timeline_manager.ppt_timeline_generated_slide import PptTimelineGeneratedSlide
from source.timeline_manager.ppt_timeline_template_slide import PptTimelineTemplateSlide

logger = logging.getLogger(__name__)


class PptTimelineFile:
    def __init__(self, input_file, output_file, timeline):
        self.output_file = output_file
        self.prs = Presentation(input_file)
        self.slides = self.prs.slides
        self.template_data = None
        self.plot_data = None
        self.template_slide = PptTimelineTemplateSlide(self.slides[0])
        self.generated_slide = PptTimelineGeneratedSlide(self.slides[1], timeline)

    def extract_template_data(self):
        self.template_data = self.template_slide.extract_template_data()

    def generate_plot_data(self):
        if self.template_data is None:
            self.extract_template_data()
        self.plot_data = self.generated_slide.generate_plot_data(self.template_data)

    def plot_timeline(self):
        if self.plot_data is None:
            self.generate_plot_data()

    def save_timeline(self):
        self.prs.save(self.output_file)

    def generate_timeline(self):
        """
        Top level method which orchestrates generation of a specific timeline.  Call other methods as required.

        :return:
        """
        self.extract_template_data()
        self.generate_plot_data()
        self.plot_timeline()
        self.save_timeline()


