import logging

from pptx.dml.color import _NoneColor, RGBColor
from pptx.util import Pt

logger = logging.getLogger(__name__)


class TemplateShape:
    """
    Wraps a shape from the template slide and use it to create new versions of the shape.  Required as there isn't a
    reliable way of cloning or copying a shape within python_pptx.
    """
    
    def __init__(self, template_shape):
        """
        Hold the source shape and capture its properties so that the shape can be used to generate new copies of itself.
        :param template_shape: 
        """
        
        self.template_shape = template_shape
        self.shape_attributes = None

    def extract_shape_attributes(self):
        """
        Takes an existing (template) autoshape and extracts formatting properties which can later be applied to
        another shape (with modifications if required).

        Note this is very simplified by making assumptions about what is likely to be needed. So for example, the fill
        colour is assumed to always be a solid colour using RGB values (for now).
        """

        fill = self.template_shape.fill
        line = self.template_shape.line
        fore_colour = fill.fore_color
        line_colour = self.template_shape.line.color
        logger.debug('line_colour = {}'.format(line_colour))
        if isinstance(line_colour._color, _NoneColor):
            line_colour_rgb = None
        else:
            line_colour_rgb = line_colour.rgb
        run = self.template_shape.text_frame.paragraphs[0].runs[0]

        self.shape_attributes = {
            'shape_type': self.template_shape.shape_type,
            'autoshape_type': self.template_shape.auto_shape_type,
            'rotation': self.template_shape.rotation,
            'left': self.template_shape.left,
            'top': self.template_shape.top,
            'width': self.template_shape.width,
            'height': self.template_shape.height,
            'fill_colour_rgb': fill.fore_color.rgb,
            'line_colour_rgb': line_colour_rgb,
            # 'line_colour_brightness': line.color.brightness,
            'line_width': line.width,
            'font_name': run.font.name,
            'font_size': run.font.size,
            'font_bold': run.font.bold,
            'font_italic': run.font.italic,
            'font_colour_rgb': run.font.color.rgb
        }

    def apply_attribute(self, shape_object, attribute_name):
        if attribute_name in self.shape_attributes and self.shape_attributes[attribute_name] is not None:
            setattr(shape_object, attribute_name, self.shape_attributes[attribute_name])
        else:
            setattr(shape_object, attribute_name, None)
    
    def apply_shape_attributes(self, target_shape, include_construction_attrs=False):
        """
        Takes a supplied (typically newly created shape) and apply this templates properties to it
    
        It's all done pretty mindlessly with each attribute processed individually.
        """
        logger.debug(
            'func:apply_shape_attributes, shape: {}, attributes: {}'.format(target_shape, self.shape_attributes))
    
        # Usually don't need width etc as that would have been set at creation of shape.
        if include_construction_attrs:
            self.apply_attribute(target_shape, 'width')
            self.apply_attribute(target_shape, 'left')
            self.apply_attribute(target_shape, 'top')
            self.apply_attribute(target_shape, 'height')
    
        if 'rotation' in self.shape_attributes and self.shape_attributes['rotation'] is not None:
            target_shape.rotation = self.shape_attributes['rotation']
    
        if 'fill_colour_rgb' in self.shape_attributes and self.shape_attributes['fill_colour_rgb'] is not None:
            fill = target_shape.fill
            fill.solid()
            fore_color = fill.fore_color
            fore_color.rgb = self.shape_attributes['fill_colour_rgb']
    
        if 'line_colour_rgb' in self.shape_attributes:
            line = target_shape.line
            if self.shape_attributes['line_colour_rgb'] is None:
                line.fill.background()
            else:
                line.color.rgb = self.shape_attributes['line_colour_rgb']
    
        if 'line_colour_brightness' in self.shape_attributes and self.shape_attributes['line_colour_brightness'] is not None:
            line = target_shape.line
            line.color.brightness = self.shape_attributes['line_colour_brightness']
    
        if 'line_width' in self.shape_attributes and self.shape_attributes['line_width'] is not None:
            line = target_shape.line
            line.width = self.shape_attributes['line_width']
    
        # Note: Only handles shapes with one text run.
        run = target_shape.text_frame.paragraphs[0].runs[0]
    
        if 'font_name' in self.shape_attributes and self.shape_attributes['font_name'] is not None:
            run.font.name = self.shape_attributes['font_name']
    
        if 'font_size' in self.shape_attributes and self.shape_attributes['font_size'] is not None:
            run.font.size = self.shape_attributes['font_size']
    
        if 'font_bold' in self.shape_attributes and self.shape_attributes['font_bold'] is not None:
            run.font.bold = self.shape_attributes['font_bold']
    
        if 'font_italic' in self.shape_attributes and self.shape_attributes['font_italic'] is not None:
            run.font.italic = self.shape_attributes['font_italic']
    
        if 'font_colour_rgb' in self.shape_attributes and self.shape_attributes['font_colour_rgb'] is not None:
            run.font.color.rgb = self.shape_attributes['font_colour_rgb']

    def create_shape_from_template(self, shapes_object, left, top, text):
        """
        Adds an autoshape to a supplied python-pptx shapes object.  The attributes of the shape are supplied and
        will have been previously extracted from a template shape which sets the formatting for all similar shapes.
    
        The top and left attributes are supplied separately as each new shape will have its own position.
    
        """
        arguments = (self.shape_attributes['autoshape_type'],
                     left,
                     top,
                     self.shape_attributes['width'],
                     self.shape_attributes['height'],
                     )
        new_shape = shapes_object.add_shape(*arguments)
    
        # Add text here, not in attributes
        text_frame = new_shape.text_frame
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.line_spacing = 0.8
        run = p.add_run()
        run.text = text
    
        # Hard code text formatting for now until we have implemented capturing these attributes
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(8)
        font.bold = True
        font.italic = False  # cause value to be inherited from theme
        font.color.rgb = RGBColor(0xFF, 0, 0)
    
        # Now apply remainder of formatting to shape
        self.apply_shape_attributes(new_shape, self.shape_attributes)