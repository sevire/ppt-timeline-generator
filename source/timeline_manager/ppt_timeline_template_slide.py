from pptx.dml.color import _NoneColor

import logging

logger = logging.getLogger(__name__)


class PptTimelineTemplateSlide:
    """
    Wraps a slide object and analyses it to find the template shapes within.  These will then be used to generate
    the shapes for the timeline.

    """
    def __init__(self, ppt_slide):
        self.template_slide = ppt_slide
        self.template_shapes = self.template_slide.shapes

    def extract_template_data(self):
        expected_templates = {'1': {}, '2': {}, '3': {}, 'TEXT 1': {}, 'TEXT 2': {}, 'TEXT 3': {}}

        for shape_number, template_shape in enumerate(self.template_shapes):
            logger.debug('{:3}: Shape type: {}'.format(shape_number, template_shape.shape_type))

            if template_shape.has_text_frame:
                shape_text = template_shape.text

                # Now work out which template shape it is and extract and store the key formatting information
                # Note: I am taking this approach as at the time of writing there is no way of simply copying a shape
                # using python-pptx
                if shape_text in expected_templates:
                    if 'data' in expected_templates[shape_text]:
                        logger.warning('Duplicate template found for {}, ignoring'.format(shape_text))
                    else:
                        # We've found a valid template so store it's data
                        logger.info('Found milestone shape {}'.format(shape_text))
                        expected_templates[shape_text] = {}
                        expected_templates[shape_text]['data'] = self.extract_shape_attributes(template_shape)
                else:
                    logger.info('Not a template')
            else:
                logger.info('No text frame')

        # We should have all templates now, just tidy up a bit and tell the user
        undiscovered = list(filter(lambda x: 'data' not in expected_templates[x], expected_templates))
        if len(undiscovered) > 0:
            logger.warning('templates not found: {}'.format(undiscovered))
        logger.info(expected_templates)

        return expected_templates

    @staticmethod
    def extract_shape_attributes(src_shape):
        """
        Takes an existing (template) autoshape and extracts formatting properties which can later be applied to
        another shape (with modifications if required).

        Note this is very simplified by making assumptions about what is likely to be needed. So for example, the fill
        colour is assumed to always be a solid colour using RGB values (for now).
        """

        fill = src_shape.fill
        line = src_shape.line
        fore_colour = fill.fore_color
        line_colour = src_shape.line.color
        logger.debug('line_colour = {}'.format(line_colour))
        if isinstance(line_colour._color, _NoneColor):
            line_colour_rgb = None
        else:
            line_colour_rgb = line_colour.rgb
        run = src_shape.text_frame.paragraphs[0].runs[0]

        shape_attributes = {
            'shape_type': src_shape.shape_type,
            'autoshape_type': src_shape.auto_shape_type,
            'rotation': src_shape.rotation,
            'left': src_shape.left,
            'top': src_shape.top,
            'width': src_shape.width,
            'height': src_shape.height,
            'fill_colour_rgb': fill.fore_color.rgb,
            'line_colour_rgb': line_colour_rgb,
            # 'line_colour_brightness': line.color.brightness,
            'line_width': line.width,
            'font_name': run.font.name,
            'font_size': run.font.size,
            'font_bold': run.font.bold,
            'font_italic': run.font.italic,
            'font_colour_rgb': run.font.color.rgb
        }

        return shape_attributes
