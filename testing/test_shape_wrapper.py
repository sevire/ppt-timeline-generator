from unittest import TestCase

from pptx import Presentation

from shape_wrapper import ShapeWrapper


class TestShapeWrapper(TestCase):
    def test_shape_wrapper(self):
        presentation_name_in = '/Users/thomasgaylardou/OneDrive - The Open University/BusinessDesign/Vision/MediumLongTerm/ACLA/EPA/timeline/testing/CriticalPathTimeline-testing-01.pptx'
        presentation_name_out = '/Users/thomasgaylardou/OneDrive - The Open University/BusinessDesign/Vision/MediumLongTerm/ACLA/EPA/timeline/testing/CriticalPathTimeline-testing-01-(out).pptx'

        prs = Presentation(presentation_name_in)

        slides = prs.slides
        template_slide = slides[0]
        slide_to_add = slides[1]

        shapes = template_slide.shapes
        new_shapes = slide_to_add.shapes

        for shape_number, template_shape in enumerate(shapes):
            print('Testing shape number {}'.format(shape_number))

            wrapped_shape = ShapeWrapper(template_shape)
            print('Shape details...\n{}'.format(wrapped_shape.__str__()))
            wrapped_shape.create_new_shape(new_shapes)

        prs.save(presentation_name_out)
